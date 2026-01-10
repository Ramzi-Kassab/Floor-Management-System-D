"""
Halliburton PowerPoint Generator - Matching Full Page View Exactly
Generates PPT files matching the web Full Page preview format.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from typing import Dict, List, Optional
import os
import io
import base64


class HalliburtonPPTGenerator:
    """Generates Halliburton-style PowerPoint presentations matching Full Page view exactly"""

    # Colors matching Full Page CSS exactly
    HALLIBURTON_RED = RGBColor(0xC8, 0x10, 0x2E)
    DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x00, 0x00, 0x00)

    # BOM Table
    TABLE_HEADER_BG = RGBColor(0x55, 0x55, 0x55)
    TABLE_ALT_ROW = RGBColor(0xF5, 0xF5, 0xF5)
    TABLE_BORDER = RGBColor(0xCC, 0xCC, 0xCC)

    # Blade Layout - matching CSS exactly
    BLADE_LABEL_BG = RGBColor(0x80, 0x80, 0x80)  # .blade-label #808080
    ROW_LABEL_BG = RGBColor(0xD0, 0xD0, 0xD0)    # .row-label #d0d0d0

    # Group colors matching CSS group-1 through group-7
    GROUP_COLORS = {
        1: (RGBColor(0xF0, 0xF0, 0xF0), RGBColor(0x33, 0x33, 0x33)),  # light bg, dark text
        2: (RGBColor(0xE0, 0xE0, 0xE0), RGBColor(0x33, 0x33, 0x33)),
        3: (RGBColor(0x80, 0x80, 0x80), RGBColor(0xFF, 0xFF, 0xFF)),  # dark bg, white text
        4: (RGBColor(0xA0, 0xA0, 0xA0), RGBColor(0xFF, 0xFF, 0xFF)),
        5: (RGBColor(0x90, 0x90, 0x90), RGBColor(0xFF, 0xFF, 0xFF)),
        6: (RGBColor(0x70, 0x70, 0x70), RGBColor(0xFF, 0xFF, 0xFF)),
        7: (RGBColor(0x60, 0x60, 0x60), RGBColor(0xFF, 0xFF, 0xFF)),
    }

    def __init__(self, page_size='Tabloid'):
        """
        Initialize with page size matching original PDF.
        
        Original PDF measurements (1283277_B.pdf):
        - Page: 792 x 1224 pt = 11" x 17" (Tabloid)
        - Cutter circle: 30.8 pt = 0.428"
        - Cutter spacing: 40.2 pt center-to-center = 0.558"
        - Drill bit image: 117.2 x 117.2 pt = 1.63"
        - Max ~17 cutters at default size, shrink if more
        """
        self.prs = Presentation()
        
        # Page size configurations based on exact PDF measurements
        if page_size == 'A4':
            self.page_width = Inches(8.27)
            self.page_height = Inches(11.69)
            self.scale = 0.70  # Calculated from 595/792 ratio
        elif page_size == 'Letter':
            self.page_width = Inches(8.5)
            self.page_height = Inches(11)
            self.scale = 0.77
        else:  # Tabloid - matches original PDF exactly
            self.page_width = Inches(11)
            self.page_height = Inches(17)
            self.scale = 1.0
        
        self.prs.slide_width = self.page_width
        self.prs.slide_height = self.page_height

        # Margins from original PDF: 24 pt = 0.33"
        self.margin_left = Inches(0.33 * self.scale)
        self.margin_top = Inches(0.24 * self.scale)
        self.content_width = self.page_width - (self.margin_left * 2)

        # Blade layout from original PDF measurements
        self.blade_label_width = Inches(0.52 * self.scale)   # 37.5 pt
        self.row_label_width = Inches(0.35 * self.scale)     # 25 pt
        
        # Cutter dimensions from original - CRITICAL for layout
        self.default_cutter_size = Inches(0.428 * self.scale)    # 30.8 pt diameter
        self.default_cutter_spacing = Inches(0.558 * self.scale)  # 40.2 pt center-to-center
        self.max_default_cutters = 17  # Max cutters before shrinking
        
        # Current values (may be adjusted dynamically per row)
        self.cutter_width = self.default_cutter_spacing
        self.cutter_circle_size = self.default_cutter_size
        
        # Row height to fit cutter + type + chamfer text
        self.row_height = Inches(0.88 * self.scale)  # ~63.5 pt per row

        # Font sizes from original PDF (in points)
        # Blade label: 33.5 pt, Row label: 16.7 pt, Position: 9.4 pt
        # Cutter type/chamfer: 5.4 pt, Group number: ~12 pt, BOM: 7.4 pt
        self.font_header_label = max(7, int(8.7 * self.scale))
        self.font_header_value = max(7, int(8.7 * self.scale))
        self.font_bom_header = max(6, int(7.4 * self.scale))
        self.font_bom_data = max(6, int(7.4 * self.scale))
        self.font_blade_label = max(20, int(33.5 * self.scale))    # B1, B2, etc
        self.font_row_label = max(12, int(16.7 * self.scale))      # R1, R2
        self.font_pos_label = max(7, int(9.4 * self.scale))        # CONE, NOSE, etc
        self.font_cutter_type = max(5, int(5.4 * self.scale))      # CT418, CT170
        self.font_cutter_group = max(10, int(12 * self.scale))     # 1, 2, 3 inside circle
        self.font_cutter_chamfer = max(5, int(5.4 * self.scale))   # 18C-60, DROP-IN
        
        # Image sizes from original PDF
        self.drill_bit_size = Inches(1.63 * self.scale)      # 117.2 pt
        self.group_shape_size = Inches(0.37 * self.scale)    # 26.8 pt
    
    def calculate_cutter_dimensions(self, num_cutters):
        """
        Dynamically calculate cutter size based on count.
        Original: ~17 cutters max at default size, shrink if more.
        """
        if num_cutters <= self.max_default_cutters:
            return self.default_cutter_size, self.default_cutter_spacing
        
        # Calculate available width for cutters
        # Available = content_width - blade_label - row_label - position_labels_space
        available = self.content_width - self.blade_label_width - self.row_label_width - Inches(0.5 * self.scale)
        
        # Shrink proportionally
        shrink_factor = float(available) / (num_cutters * float(self.default_cutter_spacing))
        if shrink_factor > 1:
            shrink_factor = 1
        
        new_size = Inches(0.428 * self.scale * shrink_factor)
        new_spacing = Inches(0.558 * self.scale * shrink_factor)
        
        return new_size, new_spacing

    @classmethod
    def get_group_color(cls, group: int):
        """Get background and text color for a group number"""
        if group in cls.GROUP_COLORS:
            return cls.GROUP_COLORS[group]
        idx = ((group - 1) % 7) + 1
        return cls.GROUP_COLORS[idx]

    def generate(self, data: dict, output_path: str) -> str:
        """Generate PowerPoint from data dictionary"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # Set white background
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = self.WHITE

        y_cursor = self.margin_top

        # Get data components
        header = data.get('header', {})
        summary = data.get('summary', [])
        blades = data.get('blades', [])
        images = data.get('images', {})
        groups = data.get('groups', [])

        # Draw header frame
        y_cursor = self._draw_header_frame(slide, header, images, y_cursor)

        # Draw content row: BOM + Group Legend + Drill bit
        y_cursor = self._draw_content_row(slide, summary, images, groups, y_cursor)

        # Small gap before blades - scaled
        y_cursor += Inches(0.15 * self.scale)

        # Draw blades
        for i, blade in enumerate(blades):
            blade_height = self._draw_blade(slide, blade, y_cursor)
            y_cursor += blade_height

            if i < len(blades) - 1:
                y_cursor += Inches(0.04 * self.scale)

            # Check if we need a new slide - based on page height
            if y_cursor > (self.page_height - Inches(0.5)) and i < len(blades) - 1:
                slide = self.prs.slides.add_slide(blank_layout)
                background = slide.background
                background.fill.solid()
                background.fill.fore_color.rgb = self.WHITE
                y_cursor = self.margin_top

        self.prs.save(output_path)
        return output_path

    def _draw_header_frame(self, slide, header: dict, images: dict, y: Emu) -> Emu:
        """Draw header frame: Logo (40% left) + Info (60% right)"""
        x = self.margin_left
        frame_width = self.content_width
        frame_height = Inches(0.45 * self.scale)

        # Frame border (rectangle with border)
        frame_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, frame_width, frame_height
        )
        frame_border.fill.background()
        frame_border.line.color.rgb = self.DARK_GRAY
        frame_border.line.width = Pt(1)

        # Logo area (40% of width)
        logo_x = x + Inches(0.1 * self.scale)
        logo_width = frame_width * 0.40

        # Try to add logo image
        logo_data = images.get('halliburton_logo', {})
        if logo_data and logo_data.get('data'):
            try:
                img_stream = self._decode_image(logo_data['data'])
                if img_stream:
                    slide.shapes.add_picture(img_stream, logo_x, y + Inches(0.1 * self.scale),
                                            width=Inches(1.0 * self.scale), height=Inches(0.25 * self.scale))
                else:
                    self._add_text_box(slide, "HALLIBURTON", logo_x, y + Inches(0.15 * self.scale),
                                      Inches(1.5 * self.scale), Inches(0.2 * self.scale), 
                                      max(10, int(14 * self.scale)),
                                      bold=True, font_color=self.HALLIBURTON_RED)
            except:
                self._add_text_box(slide, "HALLIBURTON", logo_x, y + Inches(0.15 * self.scale),
                                  Inches(1.5 * self.scale), Inches(0.2 * self.scale), 
                                  max(10, int(14 * self.scale)),
                                  bold=True, font_color=self.HALLIBURTON_RED)
        else:
            self._add_text_box(slide, "HALLIBURTON", logo_x, y + Inches(0.15 * self.scale),
                              Inches(1.5 * self.scale), Inches(0.2 * self.scale), 
                              max(10, int(14 * self.scale)),
                              bold=True, font_color=self.HALLIBURTON_RED)

        # Header info area (60% of width, starting at 40%)
        info_x = x + logo_width + Inches(0.2 * self.scale)

        # Row 1: SN Number, Mat Number, Date Created
        row1_y = y + Inches(0.08 * self.scale)
        self._add_header_item(slide, "SN Number:", header.get('sn_number', '0'),
                             info_x, row1_y, Inches(1.0 * self.scale))
        self._add_header_item(slide, "Mat Number:", header.get('mat_number', ''),
                             info_x + Inches(1.2 * self.scale), row1_y, Inches(1.3 * self.scale))
        self._add_header_item(slide, "Date Created:", header.get('date_created', ''),
                             info_x + Inches(2.6 * self.scale), row1_y, Inches(1.5 * self.scale))

        # Row 2: Revision Level, Software Version
        row2_y = y + Inches(0.25 * self.scale)
        rev = header.get('revision_level', '')
        if rev and not str(rev).startswith('D'):
            rev = f"D - {rev}"
        self._add_header_item(slide, "Revision Level:", rev,
                             info_x, row2_y, Inches(1.5 * self.scale))
        self._add_header_item(slide, "ADesc Software Version:", header.get('software_version', ''),
                             info_x + Inches(1.6 * self.scale), row2_y, Inches(2.5 * self.scale))

        return y + frame_height + Inches(0.05 * self.scale)

    def _add_header_item(self, slide, label: str, value: str, x: Emu, y: Emu, width: Emu):
        """Add a header item with label and value"""
        self._add_text_box(slide, f"{label} ", x, y, Inches(0.8 * self.scale), Inches(0.15 * self.scale),
                          self.font_header_label, font_color=self.DARK_GRAY)
        self._add_text_box(slide, str(value), x + Inches(0.5 * self.scale), y, width, Inches(0.15 * self.scale),
                          self.font_header_value, bold=True, font_color=self.DARK_GRAY)

    def _draw_content_row(self, slide, summary: list, images: dict, groups: list, y: Emu) -> Emu:
        """
        Draw content row: BOM table + Group legend + Drill bit preview
        
        Original PDF positioning:
        - Drill bit: x=646.7 pt, y=45.5 pt, size=117.2 x 117.2 pt (1.63")
        - Group shape: x=539.6 pt, size=26.8 x 24.8 pt
        """
        x = self.margin_left

        # BOM Table
        bom_width = self._draw_bom_table(slide, summary, x, y)

        # Group Legend (to the right of BOM)
        if groups:
            legend_x = x + bom_width + Inches(0.2 * self.scale)
            self._draw_group_legend(slide, groups, images, legend_x, y)

        # Drill bit face image (far right) - original position: x=646.7 pt on 792 pt page
        # That's 81.7% from left edge
        drill_data = images.get('drill_bit_face', {})
        if not (drill_data and drill_data.get('data')):
            drill_data = images.get('drill_bit_preview', {})

        if drill_data and drill_data.get('data'):
            try:
                img_stream = self._decode_image(drill_data['data'])
                if img_stream:
                    # Position at right side: page_width - margin - drill_size
                    drill_x = self.page_width - self.margin_left - self.drill_bit_size
                    slide.shapes.add_picture(img_stream, drill_x, y,
                                            width=self.drill_bit_size, height=self.drill_bit_size)
            except Exception as e:
                print(f"Failed to add drill bit image: {e}")

        # Return height based on drill bit size (largest element) + margin
        return y + max(self.drill_bit_size, Inches(1.0 * self.scale)) + Inches(0.15 * self.scale)

    def _draw_bom_table(self, slide, summary: list, x: Emu, y: Emu) -> Emu:
        """Draw BOM table and return width used"""
        if not summary:
            return Inches(0)

        # Column widths - scaled dynamically
        col_widths = [
            Inches(0.22 * self.scale),   # #
            Inches(0.45 * self.scale),   # Size
            Inches(0.55 * self.scale),   # Chamfer
            Inches(0.55 * self.scale),   # Type
            Inches(0.35 * self.scale),   # Count
            Inches(0.7 * self.scale),    # Mat #
            Inches(0.65 * self.scale)    # Family #
        ]
        headers = ["#", "Size", "Chamfer", "Type", "Count", "Mat #", "Family #"]

        total_width = sum(col_widths, Inches(0))
        row_height = Inches(0.16 * self.scale)
        header_height = Inches(0.18 * self.scale)

        # Header row background
        header_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, total_width, header_height)
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = self.TABLE_HEADER_BG
        header_bg.line.color.rgb = self.DARK_GRAY
        header_bg.line.width = Pt(0.5)

        # Header text
        col_x = x
        for header_text, width in zip(headers, col_widths):
            self._add_text_box(slide, header_text, col_x, y + Inches(0.02 * self.scale), width, Inches(0.12 * self.scale),
                              self.font_bom_header, bold=True, font_color=self.WHITE, align='center')
            col_x += width

        # Data rows
        current_y = y + header_height
        for row_idx, row in enumerate(summary):
            # Alternating row colors
            row_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, current_y, total_width, row_height)
            if row_idx % 2 == 0:
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = self.TABLE_ALT_ROW
            else:
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = self.WHITE
            row_bg.line.color.rgb = self.LIGHT_GRAY
            row_bg.line.width = Pt(0.5)

            row_data = [
                str(row.get('index', '')),
                str(row.get('size', '')),
                str(row.get('chamfer', '')),
                str(row.get('type', '')),
                str(row.get('count', '')),
                str(row.get('mat_number', '')),
                str(row.get('family_number', ''))
            ]

            col_x = x
            for value, width in zip(row_data, col_widths):
                self._add_text_box(slide, value, col_x, current_y + Inches(0.01 * self.scale),
                                  width, Inches(0.12 * self.scale),
                                  self.font_bom_data, align='center')
                col_x += width

            current_y += row_height

        return total_width

    def _draw_group_legend(self, slide, groups: list, images: dict, x: Emu, y: Emu):
        """Draw group legend with shape image and group values"""
        # Try to draw shape image - scaled
        shape_data = images.get('group_shape', {})
        shape_offset = Inches(0)

        if shape_data and shape_data.get('data'):
            try:
                img_stream = self._decode_image(shape_data['data'])
                if img_stream:
                    slide.shapes.add_picture(img_stream, x, y, 
                                            width=self.group_shape_size, height=self.group_shape_size)
                    shape_offset = self.group_shape_size + Inches(0.05 * self.scale)
            except Exception as e:
                print(f"Failed to add group shape: {e}")

        # "Group" header
        self._add_text_box(slide, "Group", x + shape_offset, y + Inches(0.1 * self.scale),
                          Inches(0.5 * self.scale), Inches(0.15 * self.scale), 
                          max(6, int(8 * self.scale)), bold=True, align='center')

        # Group values
        group_text = ', '.join(str(g) for g in groups)
        self._add_text_box(slide, group_text, x + shape_offset, y + Inches(0.3 * self.scale),
                          Inches(0.5 * self.scale), Inches(0.2 * self.scale), 
                          max(8, int(10 * self.scale)), bold=True, align='center')

    def _draw_blade(self, slide, blade: dict, y_top: Emu) -> Emu:
        """Draw a single blade matching Full Page layout exactly"""
        blade_name = blade.get('name', 'B1')
        x = self.margin_left

        # Get rows with data
        rows_to_draw = []
        for row_key in ['r1', 'r2', 'r3', 'r4']:
            row_data = blade.get(row_key, {})
            if row_data and self._has_row_data(row_data):
                rows_to_draw.append((row_key.upper(), row_data))

        if not rows_to_draw:
            return Inches(0)

        num_rows = len(rows_to_draw)
        total_height = self.row_height * num_rows

        # Blade label background (gray, spans all rows)
        blade_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y_top,
            self.blade_label_width, total_height
        )
        blade_bg.fill.solid()
        blade_bg.fill.fore_color.rgb = self.BLADE_LABEL_BG
        blade_bg.line.fill.background()

        # Blade name (centered vertically)
        label_y = y_top + total_height / 2 - Inches(0.08 * self.scale)
        self._add_text_box(slide, blade_name, x, label_y,
                          self.blade_label_width, Inches(0.16 * self.scale),
                          self.font_blade_label, bold=True, font_color=self.WHITE, align='center')

        # Draw each row
        y_current = y_top
        row_x = x + self.blade_label_width

        for row_name, row_data in rows_to_draw:
            self._draw_blade_row(slide, row_name, row_data, row_x, y_current)
            y_current += self.row_height

        return total_height + Inches(0.03 * self.scale)

    def _draw_blade_row(self, slide, row_name: str, row_data: dict, x: Emu, y: Emu):
        """
        Draw a single blade row with row label and cutters.
        
        Uses dynamic sizing: if >17 cutters, they shrink to fit.
        Original PDF uses 30.8 pt circles with 40.2 pt spacing for ≤17 cutters.
        """
        # Row label background (light gray)
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y,
            self.row_label_width, self.row_height
        )
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = self.ROW_LABEL_BG
        row_bg.line.fill.background()

        # Row name (R1, R2) - font size 16.7 pt in original
        row_label_y = y + self.row_height / 2 - Inches(0.12 * self.scale)
        self._add_text_box(slide, row_name, x, row_label_y,
                          self.row_label_width, Inches(0.24 * self.scale),
                          self.font_row_label, bold=True, align='center')

        # Count total cutters in this row for dynamic sizing
        positions_order = ["CONE", "NOSE", "SHOULDER", "GAUGE", "PAD"]
        total_cutters = sum(len(row_data.get(pos, [])) for pos in positions_order)
        num_positions = sum(1 for pos in positions_order if row_data.get(pos))
        
        # Calculate cutter dimensions (dynamic if >17)
        cutter_size, cutter_spacing = self.calculate_cutter_dimensions(total_cutters)
        
        # Position label width - scaled based on cutter size
        pos_label_width = Inches(0.15 * self.scale)
        
        # Start drawing cutters
        cutter_x = x + self.row_label_width + Inches(0.05 * self.scale)

        for pos in positions_order:
            cells = row_data.get(pos, [])
            if cells:
                # Position label (CONE, NOSE, etc.) - rotated vertical text
                # Font size 9.4 pt in original
                pos_label_y = y + Inches(0.05 * self.scale)
                self._add_text_box(slide, pos, cutter_x, pos_label_y,
                                  pos_label_width, self.row_height - Inches(0.1 * self.scale),
                                  self.font_pos_label, bold=True, align='center')
                cutter_x += pos_label_width

                # Draw each cutter with dynamic sizing
                for cell in cells:
                    self._draw_cutter_dynamic(slide, cell, cutter_x, y, cutter_size, cutter_spacing)
                    cutter_x += cutter_spacing
    
    def _draw_cutter_dynamic(self, slide, cell: dict, x: Emu, y: Emu, circle_size: Emu, spacing: Emu):
        """
        Draw a single cutter with dynamic size.
        
        Layout from original PDF:
        - Type text (CT418, CT170): 5.4 pt, above circle
        - Circle: 30.8 pt diameter (or scaled)
        - Group number: 12 pt, inside circle
        - Chamfer text (18C-60): 5.4 pt, below circle
        """
        cutter_type = cell.get('type', '')
        group = cell.get('group', 1)
        chamfer = cell.get('chamfer', '')

        # Center the circle within the spacing
        center_x = x + spacing / 2 - circle_size / 2

        # Cutter type (top) - 5.4 pt font
        type_y = y + Inches(0.03 * self.scale)
        self._add_text_box(slide, cutter_type, x, type_y,
                          spacing, Inches(0.15 * self.scale),
                          self.font_cutter_type, align='center')

        # Circle with group number - positioned in middle of row height
        circle_y = y + Inches(0.20 * self.scale)
        bg_color, text_color = self.get_group_color(group)

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, center_x, circle_y,
            circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = bg_color
        circle.line.color.rgb = self.DARK_GRAY
        circle.line.width = Pt(1)

        # Group number in circle - 12 pt font
        group_text_y = circle_y + circle_size / 2 - Inches(0.08 * self.scale)
        self._add_text_box(slide, str(group), center_x, group_text_y,
                          circle_size, Inches(0.16 * self.scale),
                          self.font_cutter_group, bold=True, font_color=text_color, align='center')

        # Chamfer (bottom) - 5.4 pt font
        if chamfer:
            chamfer_y = y + self.row_height - Inches(0.18 * self.scale)
            self._add_text_box(slide, chamfer, x, chamfer_y,
                              spacing, Inches(0.15 * self.scale),
                              self.font_cutter_chamfer, align='center')

    def _add_text_box(self, slide, text: str, x: Emu, y: Emu,
                     width: Emu, height: Emu,
                     font_size: int = 10, bold: bool = False,
                     font_color: RGBColor = None, align: str = 'left'):
        """Add a text box to the slide"""
        if font_color is None:
            font_color = self.BLACK

        textbox = slide.shapes.add_textbox(x, y, width, height)
        tf = textbox.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.margin_left = Pt(1)
        tf.margin_right = Pt(1)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)

        p = tf.paragraphs[0]
        p.text = str(text) if text else ''
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.font.name = 'Arial'

        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

    def _has_row_data(self, row_data: dict) -> bool:
        """Check if row has any cutter data"""
        if not row_data:
            return False
        for pos in ["CONE", "NOSE", "SHOULDER", "GAUGE", "PAD"]:
            if row_data.get(pos):
                return True
        return False

    def _decode_image(self, data_url: str):
        """Decode base64 data URL to BytesIO stream"""
        if not data_url or not data_url.startswith('data:'):
            return None

        try:
            header, encoded = data_url.split(',', 1)
            img_data = base64.b64decode(encoded)
            return io.BytesIO(img_data)
        except:
            return None


def generate_ppt_from_pdf(pdf_path: str, output_path: str = None) -> str:
    """Extract data from PDF and generate matching PPT."""
    from pdf_extractor_fitz import extract_pdf_data

    data = extract_pdf_data(pdf_path)

    if output_path is None:
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        output_dir = os.path.dirname(pdf_path)
        output_path = os.path.join(output_dir, f"{base_name}.pptx")

    generator = HalliburtonPPTGenerator()
    generator.generate(data, output_path)

    return output_path


def batch_generate_ppt(input_folder: str, output_folder: str = None) -> List[dict]:
    """Generate PPT files for all PDFs in a folder."""
    if output_folder is None:
        output_folder = os.path.join(os.path.dirname(input_folder), 'outputs')

    os.makedirs(output_folder, exist_ok=True)

    results = []
    pdf_files = [f for f in os.listdir(input_folder) if f.lower().endswith('.pdf')]

    print(f"Found {len(pdf_files)} PDF files to process")

    for pdf_file in sorted(pdf_files):
        pdf_path = os.path.join(input_folder, pdf_file)
        ppt_file = os.path.splitext(pdf_file)[0] + '.pptx'
        ppt_path = os.path.join(output_folder, ppt_file)

        result = {'input': pdf_path, 'output': ppt_path, 'success': False, 'error': None}

        try:
            generate_ppt_from_pdf(pdf_path, ppt_path)
            result['success'] = True
            print(f"  OK: {pdf_file} -> {ppt_file}")
        except Exception as e:
            result['error'] = str(e)
            print(f"  FAILED: {pdf_file} - {e}")

        results.append(result)

    success_count = sum(1 for r in results if r['success'])
    print(f"\nProcessed {len(results)} files: {success_count} success, {len(results) - success_count} failed")

    return results


if __name__ == '__main__':
    import sys

    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    ORIGINALS_DIR = os.path.join(BASE_DIR, 'originals')
    OUTPUTS_DIR = os.path.join(BASE_DIR, 'outputs')

    if len(sys.argv) > 1:
        pdf_path = sys.argv[1]
        if os.path.exists(pdf_path):
            output_path = os.path.join(OUTPUTS_DIR,
                                      os.path.splitext(os.path.basename(pdf_path))[0] + '.pptx')
            result = generate_ppt_from_pdf(pdf_path, output_path)
            print(f"Generated: {result}")
        else:
            print(f"File not found: {pdf_path}")
    else:
        batch_generate_ppt(ORIGINALS_DIR, OUTPUTS_DIR)
